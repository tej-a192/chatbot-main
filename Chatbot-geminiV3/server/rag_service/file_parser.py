# server/rag_service/file_parser.py
import os
try:
    import pypdf
except ImportError:
    print("pypdf not found, PDF parsing will fail. Install with: pip install pypdf")
    pypdf = None # Set to None if not installed

try:
    from docx import Document as DocxDocument
except ImportError:
    print("python-docx not found, DOCX parsing will fail. Install with: pip install python-docx")
    DocxDocument = None

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangchainDocument
from rag_service import config # Import from package
import logging

# Configure logger for this module
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO) # Or DEBUG for more details
handler = logging.StreamHandler()
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
handler.setFormatter(formatter)
if not logger.hasHandlers():
    logger.addHandler(handler)


def parse_pdf(file_path):
    """Extracts text content from a PDF file using pypdf."""
    if not pypdf: return None # Check if library loaded
    text = ""
    try:
        reader = pypdf.PdfReader(file_path)
        num_pages = len(reader.pages)
        # logger.debug(f"Reading {num_pages} pages from PDF: {os.path.basename(file_path)}")
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n" # Add newline between pages
            except Exception as page_err:
                 logger.warning(f"Error extracting text from page {i+1} of {os.path.basename(file_path)}: {page_err}")
        # logger.debug(f"Extracted {len(text)} characters from PDF.")
        return text.strip() if text.strip() else None # Return None if empty after stripping
    except FileNotFoundError:
        logger.error(f"PDF file not found: {file_path}")
        return None
    except pypdf.errors.PdfReadError as pdf_err:
        logger.error(f"Error reading PDF {os.path.basename(file_path)} (possibly corrupted or encrypted): {pdf_err}")
        return None
    except Exception as e:
        logger.error(f"Unexpected error parsing PDF {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

def parse_docx(file_path):
    """Extracts text content from a DOCX file."""
    if not DocxDocument: return None # Check if library loaded
    try:
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        # logger.debug(f"Extracted {len(text)} characters from DOCX.")
        return text.strip() if text.strip() else None
    except Exception as e:
        logger.error(f"Error parsing DOCX {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

def parse_txt(file_path):
    """Reads text content from a TXT file (or similar plain text like .py, .js)."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        # logger.debug(f"Read {len(text)} characters from TXT file.")
        return text.strip() if text.strip() else None
    except Exception as e:
        logger.error(f"Error parsing TXT {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

# Add PPTX parsing (requires python-pptx)
try:
    from pptx import Presentation
    PPTX_SUPPORTED = True
    def parse_pptx(file_path):
        """Extracts text content from a PPTX file."""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = shape.text.strip()
                        if shape_text:
                            text += shape_text + "\n" # Add newline between shape texts
            # logger.debug(f"Extracted {len(text)} characters from PPTX.")
            return text.strip() if text.strip() else None
        except Exception as e:
            logger.error(f"Error parsing PPTX {os.path.basename(file_path)}: {e}", exc_info=True)
            return None
except ImportError:
    PPTX_SUPPORTED = False
    logger.warning("python-pptx not installed. PPTX parsing will be skipped.")
    def parse_pptx(file_path):
        logger.warning(f"Skipping PPTX file {os.path.basename(file_path)} as python-pptx is not installed.")
        return None


def parse_file(file_path):
    """Parses a file based on its extension, returning text content or None."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    logger.debug(f"Attempting to parse file: {os.path.basename(file_path)} (Extension: {ext})")

    if ext == '.pdf':
        return parse_pdf(file_path)
    elif ext == '.docx':
        return parse_docx(file_path)
    elif ext == '.pptx':
        return parse_pptx(file_path) # Use the conditional function
    elif ext in ['.txt', '.py', '.js', '.md', '.log', '.csv', '.html', '.xml', '.json']: # Expand text-like types
        return parse_txt(file_path)
    # Add other parsers here if needed (e.g., for .doc, .xls)
    elif ext == '.doc':
        # Requires antiword or similar external tool, more complex
        logger.warning(f"Parsing for legacy .doc files is not implemented: {os.path.basename(file_path)}")
        return None
    else:
        logger.warning(f"Unsupported file extension for parsing: {ext} ({os.path.basename(file_path)})")
        return None

def chunk_text(text, file_name, user_id):
    """Chunks text and creates Langchain Documents with metadata."""
    if not text or not isinstance(text, str):
        logger.warning(f"Invalid text input for chunking (file: {file_name}). Skipping.")
        return []

    # Use splitter configured in config.py
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        length_function=len,
        is_separator_regex=False, # Use default separators
        # separators=["\n\n", "\n", " ", ""] # Default separators
    )

    try:
        chunks = text_splitter.split_text(text)
        if not chunks:
             logger.warning(f"Text splitting resulted in zero chunks for file: {file_name}")
             return []

        documents = []
        for i, chunk in enumerate(chunks):
             # Ensure chunk is not just whitespace before creating Document
             if chunk and chunk.strip():
                 documents.append(
                     LangchainDocument(
                         page_content=chunk,
                         metadata={
                             'userId': user_id, # Store user ID
                             'documentName': file_name, # Store original filename
                             'chunkIndex': i # Store chunk index for reference
                         }
                     )
                 )
        if documents:
            logger.info(f"Split '{file_name}' into {len(documents)} non-empty chunks.")
        else:
            logger.warning(f"No non-empty chunks created for file: {file_name} after splitting.")
        return documents
    except Exception as e:
        logger.error(f"Error during text splitting for file {file_name}: {e}", exc_info=True)
        return [] # Return empty list on error
